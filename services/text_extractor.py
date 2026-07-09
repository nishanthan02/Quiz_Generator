# services/text_extractor.py
# ============================================================
# Text extraction and chunking logic — refactored from the
# original prototype's file-reading section.
#
# Supported formats:
#   .pdf  — extracted with PyMuPDF (fitz); preserves paragraph
#           structure better than pdfminer for lecture slides.
#   .pptx — extracted with python-pptx; iterates over every
#            text frame in every slide.
#
# Chunking strategy:
#   Fixed-size word-count chunks (default 250 words) with a
#   50-word overlap so questions about content near chunk
#   boundaries still get useful context.
#   This mirrors the original but adds overlap for RAG quality.
# ============================================================

from pptx import Presentation
import fitz  # PyMuPDF


# ── Extraction ───────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract all text from a PDF given its raw bytes.
    Uses PyMuPDF which handles modern PDFs (including those
    exported from PowerPoint) much better than pdfminer.
    """
    text_parts: list[str] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            page_text = page.get_text("text")  # plain text, no layout
            if page_text.strip():
                text_parts.append(page_text.strip())
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract all visible text from a .pptx file.
    Iterates over every slide → every shape → every text frame
    to capture all bullet points and title text.
    """
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    text_runs: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
    return "\n".join(text_runs)


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Dispatch to the correct extractor based on file extension.
    Raises ValueError for unsupported types — the Celery task
    will catch this and mark the document as 'failed'.
    """
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "pptx":
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file type '.{ext}'. "
            "Only .pdf and .pptx are accepted."
        )


# ── Chunking ─────────────────────────────────────────────────

def chunk_text(
    text: str,
    chunk_size: int = 250,
    overlap: int = 50,
) -> list[str]:
    """
    Split text into fixed-size word chunks with overlap.

    Args:
        text:       The full extracted text string.
        chunk_size: Target number of words per chunk.
        overlap:    Number of words to repeat from the previous
                    chunk to preserve cross-boundary context.

    Returns:
        List of chunk strings.

    Example with chunk_size=5, overlap=2:
        words = [A, B, C, D, E, F, G]
        chunks = ["A B C D E", "D E F G"]
    """
    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    step = chunk_size - overlap  # how far to advance each iteration

    for start in range(0, len(words), step):
        chunk = " ".join(words[start : start + chunk_size])
        if chunk.strip():
            chunks.append(chunk)

        # Stop if we've covered all words
        if start + chunk_size >= len(words):
            break

    return chunks
